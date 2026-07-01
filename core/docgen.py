# -*- coding: utf-8 -*-
"""
core/docgen.py — สร้างเอกสาร Word/Excel/PowerPoint ด้วย AI + Python
- Word: AI ร่างเนื้อหา → python-docx สร้าง .docx (หัวข้อ, ย่อหน้า, ตาราง, ฟอนต์ไทย)
- Excel: AI จัดข้อมูลเป็น JSON → openpyxl สร้าง .xlsx (สูตร, จัดรูปแบบ)
- PowerPoint: AI สร้างเนื้อหาสไลด์ → python-pptx สร้าง .pptx
"""
import os
import json
import re

_BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_EXPORTS_DIR = os.path.join(_BASE, "exports")
os.makedirs(_EXPORTS_DIR, exist_ok=True)

THAI_FONT = "TH Sarabun New"

# ──────────────────────────────────────────────
# Word (.docx)
# ──────────────────────────────────────────────

def generate_word(topic, content_hint="", username=None, llm=None):
    """AI ร่างเนื้อหา → สร้างไฟล์ .docx
    คืน (filepath, ai_note)
    """
    from docx import Document
    from docx.shared import Pt, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    # 1. AI สร้างโครงสร้างเอกสารเป็น JSON
    system_prompt = (
        "คุณเป็นผู้ช่วยเขียนเอกสาร สร้างเนื้อหาเป็น JSON ตามรูปแบบที่กำหนด\n"
        "ตอบเป็น JSON เท่านั้น ไม่ต้องมีคำอธิบายนอก JSON\n"
        f"ใช้ภาษาไทย เนื้อหาเป็นทางการ อ่านง่าย"
    )
    user_prompt = (
        f'หัวข้อเอกสาร: {topic}\n'
        f'คำแนะนำเพิ่มเติม: {content_hint or "ไม่มี"}\n\n'
        'สร้างเอกสารเป็น JSON รูปแบบนี้:\n'
        '```json\n'
        '{\n'
        '  "title": "ชื่อเอกสาร",\n'
        '  "subtitle": "บทย่อย (ถ้ามี)",\n'
        '  "author": "ผู้เขียน",\n'
        '  "sections": [\n'
        '    {\n'
        '      "heading": "หัวข้อ 1",\n'
        '      "level": 1,\n'
        '      "paragraphs": ["ย่อหน้าที่ 1", "ย่อหน้าที่ 2"],\n'
        '      "bullets": ["ข้อ 1", "ข้อ 2"],\n'
        '      "table": {"headers": ["คอลัมน์ 1", "คอลัมน์ 2"], "rows": [["ข้อมูล", "ข้อมูล"]]}\n'
        '    }\n'
        '  ]\n'
        '}\n'
        '```\n'
        'หมายเหตุ: แต่ละ section มี heading + อย่างน้อย paragraphs หรือ bullets หรือ table อย่างใดอย่างหนึ่ง\n'
        'level 1 = หัวข้อใหญ่, level 2 = หัวข้อย่อย'
    )

    if llm is None:
        from core import llm as llm
    reply, tier, _ = llm.call_tier("deepseek", system_prompt, user_prompt, max_tokens=6000)

    # 2. Parse JSON จาก AI
    doc_data = _extract_json(reply)
    if not doc_data:
        return None, f"AI ไม่ส่ง JSON ที่ถูกต้อง: {reply[:200]}"

    # 3. สร้าง .docx
    doc = Document()

    # ตั้งฟอนต์ไทยเป็น default
    style = doc.styles['Normal']
    style.font.name = THAI_FONT
    style.font.size = Pt(16)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), THAI_FONT)

    # หน้ากระดาษ A4
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(3)
    section.right_margin = Cm(2.5)

    # หัวเรื่อง
    title = doc_data.get("title", topic)
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _set_font_recursive(h, THAI_FONT, Pt(28))

    subtitle = doc_data.get("subtitle")
    if subtitle:
        p = doc.add_paragraph(subtitle)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _set_font_recursive(p, THAI_FONT, Pt(18))

    author = doc_data.get("author")
    if author:
        p = doc.add_paragraph(f"ผู้เขียน: {author}")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _set_font_recursive(p, THAI_FONT, Pt(14))

    doc.add_paragraph()  # บรรทัดว่าง

    # แต่ละ section
    for sec in doc_data.get("sections", []):
        level = sec.get("level", 1)
        heading_text = sec.get("heading", "")
        if heading_text:
            h = doc.add_heading(heading_text, level=min(level, 3))
            _set_font_recursive(h, THAI_FONT, Pt(20 if level == 1 else 16))

        for para in sec.get("paragraphs", []):
            p = doc.add_paragraph(para)
            _set_font_recursive(p, THAI_FONT, Pt(16))

        for bullet in sec.get("bullets", []):
            p = doc.add_paragraph(bullet, style='List Bullet')
            _set_font_recursive(p, THAI_FONT, Pt(16))

        tbl = sec.get("table")
        if tbl and tbl.get("headers") and tbl.get("rows"):
            headers = tbl["headers"]
            rows = tbl["rows"]
            table = doc.add_table(rows=1 + len(rows), cols=len(headers))
            table.style = 'Table Grid'
            for i, hdr in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = hdr
                _set_font_recursive(cell.paragraphs[0], THAI_FONT, Pt(14))
            for r, row in enumerate(rows):
                for c, val in enumerate(row):
                    if c < len(headers):
                        cell = table.rows[r + 1].cells[c]
                        cell.text = str(val)
                        _set_font_recursive(cell.paragraphs[0], THAI_FONT, Pt(14))

    # บันทึก
    slug = _slugify(title)
    filepath = os.path.join(_EXPORTS_DIR, f"{slug}.docx")
    doc.save(filepath)
    return filepath, f"สร้างด้วย {tier}"


# ──────────────────────────────────────────────
# Excel (.xlsx)
# ──────────────────────────────────────────────

def generate_excel(topic, data_hint="", username=None, llm=None):
    """AI จัดข้อมูลเป็น JSON → สร้างไฟล์ .xlsx
    คืน (filepath, ai_note)
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    system_prompt = (
        "คุณเป็นผู้ช่วยจัดข้อมูลเป็นตาราง Excel สร้างข้อมูลเป็น JSON ตามรูปแบบที่กำหนด\n"
        "ตอบเป็น JSON เท่านั้น ไม่ต้องมีคำอธิบายนอก JSON\n"
        f"ใช้ภาษาไทย ข้อมูลถูกต้อง ตรงประเด็น"
    )
    user_prompt = (
        f'หัวข้อตาราง: {topic}\n'
        f'ข้อมูล/คำอธิบาย: {data_hint or "ไม่มี"}\n\n'
        'สร้างตารางเป็น JSON รูปแบบนี้:\n'
        '```json\n'
        '{\n'
        '  "sheet_name": "ชื่อชีท",\n'
        '  "title": "ชื่อตาราง",\n'
        '  "headers": ["คอลัมน์ 1", "คอลัมน์ 2", "คอลัมน์ 3"],\n'
        '  "rows": [["ข้อมูล", 123, "ข้อความ"], ["ข้อมูล", 456, "ข้อความ"]],\n'
        '  "sum_columns": [1],\n'
        '  "avg_columns": [],\n'
        '  "note": "หมายเหตุถ้ามี"\n'
        '}\n'
        '```\n'
        'หมายเหตุ: sum_columns = คอลัมน์ที่ต้องการสูตร SUM (index เริ่ม 0)\n'
        'avg_columns = คอลัมน์ที่ต้องการสูตร AVERAGE\n'
        'ตัวเลขใน rows ให้เป็น number ไม่ใช่ string'
    )

    if llm is None:
        from core import llm as llm
    reply, tier, _ = llm.call_tier("deepseek", system_prompt, user_prompt, max_tokens=4000)

    tbl_data = _extract_json(reply)
    if not tbl_data:
        return None, f"AI ไม่ส่ง JSON ที่ถูกต้อง: {reply[:200]}"

    wb = Workbook()
    ws = wb.active
    ws.title = tbl_data.get("sheet_name", "Sheet1")[:31]

    # สไตล์
    header_font = Font(name=THAI_FONT, size=14, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    body_font = Font(name=THAI_FONT, size=12)
    body_align = Alignment(vertical="center", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin")
    )

    # ชื่อตาราง (row 1)
    title = tbl_data.get("title", topic)
    ws.cell(row=1, column=1, value=title)
    ws.cell(row=1, column=1).font = Font(name=THAI_FONT, size=16, bold=True)
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(tbl_data.get("headers", [])))
    ws.cell(row=1, column=1).alignment = Alignment(horizontal="center", vertical="center")

    # Headers (row 2)
    headers = tbl_data.get("headers", [])
    for c, hdr in enumerate(headers, 1):
        cell = ws.cell(row=2, column=c, value=hdr)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border

    # Data rows
    rows = tbl_data.get("rows", [])
    for r, row in enumerate(rows, 3):
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=r, column=c, value=val)
            cell.font = body_font
            cell.alignment = body_align
            cell.border = thin_border

    # สูตร SUM / AVERAGE
    sum_cols = tbl_data.get("sum_columns", [])
    avg_cols = tbl_data.get("avg_columns", [])
    last_data_row = 2 + len(rows)
    if sum_cols or avg_cols:
        formula_row = last_data_row + 1
        for c in range(1, len(headers) + 1):
            cell = ws.cell(row=formula_row, column=c)
            cell.font = Font(name=THAI_FONT, size=12, bold=True)
            cell.border = thin_border
            col_letter = get_column_letter(c)
            if (c - 1) in sum_cols:
                cell.value = f"=SUM({col_letter}3:{col_letter}{last_data_row})"
                cell.alignment = Alignment(horizontal="right")
            elif (c - 1) in avg_cols:
                cell.value = f"=AVERAGE({col_letter}3:{col_letter}{last_data_row})"
                cell.alignment = Alignment(horizontal="right")
            elif c == 1:
                cell.value = "รวม"

    # ความกว้างคอลัมน์
    for c in range(1, len(headers) + 1):
        ws.column_dimensions[get_column_letter(c)].width = 20

    # หมายเหตุ
    note = tbl_data.get("note")
    if note:
        nr = formula_row + 2 if (sum_cols or avg_cols) else last_data_row + 2
        cell = ws.cell(row=nr, column=1, value=f"หมายเหตุ: {note}")
        cell.font = Font(name=THAI_FONT, size=11, italic=True)

    # บันทึก
    slug = _slugify(title)
    filepath = os.path.join(_EXPORTS_DIR, f"{slug}.xlsx")
    wb.save(filepath)
    return filepath, f"สร้างด้วย {tier}"


# ──────────────────────────────────────────────
# PowerPoint (.pptx)
# ──────────────────────────────────────────────

def generate_pptx(topic, content_hint="", username=None, llm=None):
    """AI สร้างเนื้อหาสไลด์ → สร้างไฟล์ .pptx
    คืน (filepath, ai_note)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    system_prompt = (
        "คุณเป็นผู้ช่วยสร้างงานนำเสนอ สร้างเนื้อหาสไลด์เป็น JSON ตามรูปแบบที่กำหนด\n"
        "ตอบเป็น JSON เท่านั้น ไม่ต้องมีคำอธิบายนอก JSON\n"
        f"ใช้ภาษาไทย เนื้อหากระชับ ชัดเจน เหมาะกับการนำเสนอ"
    )
    user_prompt = (
        f'หัวข้อการนำเสนอ: {topic}\n'
        f'คำแนะนำเพิ่มเติม: {content_hint or "ไม่มี"}\n\n'
        'สร้างสไลด์เป็น JSON รูปแบบนี้:\n'
        '```json\n'
        '{\n'
        '  "title": "ชื่อการนำเสนอ",\n'
        '  "subtitle": "บทย่อย",\n'
        '  "author": "ผู้นำเสนอ",\n'
        '  "slides": [\n'
        '    {\n'
        '      "layout": "title",\n'
        '      "title": "หัวข้อสไลด์",\n'
        '      "subtitle": "บทย่อย",\n'
        '      "bullets": ["ข้อ 1", "ข้อ 2", "ข้อ 3"],\n'
        '      "notes": "บันทึกผู้พูด"\n'
        '    },\n'
        '    {\n'
        '      "layout": "content",\n'
        '      "title": "หัวข้อเนื้อหา",\n'
        '      "bullets": ["ข้อ 1", "ข้อ 2"],\n'
        '      "table": {"headers": ["col1", "col2"], "rows": [["a", "b"]]},\n'
        '      "notes": ""\n'
        '    }\n'
        '  ]\n'
        '}\n'
        '```\n'
        'layout มี: "title" (สไลด์ปก), "content" (เนื้อหา), "section" (แบ่งบท), "closing" (สรุป/ขอบคุณ)\n'
        'สไลด์แรกใช้ layout "title", สไลด์สุดท้ายใช้ layout "closing"\n'
        'bullets ไม่เกิน 5 ข้อต่อสไลด์ แต่ละข้อไม่เกิน 2 บรรทัด'
    )

    if llm is None:
        from core import llm as llm
    reply, tier, _ = llm.call_tier("deepseek", system_prompt, user_prompt, max_tokens=6000)

    pptx_data = _extract_json(reply)
    if not pptx_data:
        return None, f"AI ไม่ส่ง JSON ที่ถูกต้อง: {reply[:200]}"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_color = RGBColor(0x2C, 0x3E, 0x50)
    accent_color = RGBColor(0x4A, 0x6F, 0xA5)
    bg_color = RGBColor(0xF8, 0xF9, 0xFA)

    for slide_data in pptx_data.get("slides", []):
        layout_type = slide_data.get("layout", "content")
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # พื้นหลัง
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color

        if layout_type == "title":
            _add_title_slide(slide, slide_data, title_color, accent_color, THAI_FONT)
        elif layout_type == "section":
            _add_section_slide(slide, slide_data, title_color, accent_color, THAI_FONT)
        elif layout_type == "closing":
            _add_closing_slide(slide, slide_data, title_color, accent_color, THAI_FONT)
        else:
            _add_content_slide(slide, slide_data, title_color, accent_color, THAI_FONT)

        # notes
        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    slug = _slugify(pptx_data.get("title", topic))
    filepath = os.path.join(_EXPORTS_DIR, f"{slug}.pptx")
    prs.save(filepath)
    return filepath, f"สร้างด้วย {tier}"


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _extract_json(text):
    """สกัด JSON จาก text ที่อาจมี ```json block"""
    # ลองหา ```json ... ```
    m = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', text, re.DOTALL)
    if m:
        try:
            return json.loads(m.group(1))
        except json.JSONDecodeError:
            pass
    # ลอง parse ตรง
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError:
        pass
    # หา { ... } แรก
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end != -1:
        try:
            return json.loads(text[start:end+1])
        except json.JSONDecodeError:
            pass
    return None


def _slugify(text):
    """สร้างชื่อไฟล์ที่ปลอดภัยจากข้อความไทย"""
    import unicodedata
    # เก็บตัวอักษรไทย + อังกฤษ + ตัวเลข
    safe = "".join(c for c in text if c.isalnum() or 0x0e01 <= ord(c) <= 0x0e5b or c in '-_ ')
    safe = safe.strip().replace(' ', '_')[:50]
    if not safe:
        safe = "document"
    return safe


def _set_font_recursive(paragraph, font_name, font_size):
    """ตั้งฟอนต์ไทยให้ทุก run ใน paragraph"""
    from docx.oxml.ns import qn
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = font_size
        run.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)


def _add_textbox(slide, left, top, width, height):
    """เพิ่ม textbox แล้วคืน text_frame"""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _set_pptx_font(run, font_name, font_size, color=None, bold=False):
    """ตั้งฟอนต์สำหรับ run ใน pptx"""
    from pptx.util import Pt
    run.font.name = font_name
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    run.font.bold = bold


def _add_title_slide(slide, data, title_color, accent_color, font_name):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # แถบสีบน
    from pptx.shapes.autoshape import Shape
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(2.5))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    shape.line.fill.background()

    tf = _add_textbox(slide, 1, 0.6, 11, 1.5)
    tf.text = data.get("title", "")
    _set_pptx_font(tf.paragraphs[0].runs[0], font_name, 36, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = data.get("subtitle", "")
    if subtitle:
        tf2 = _add_textbox(slide, 1, 2.8, 11, 1)
        tf2.text = subtitle
        _set_pptx_font(tf2.paragraphs[0].runs[0], font_name, 22, accent_color)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    author = data.get("author", "")
    if author:
        tf3 = _add_textbox(slide, 1, 5.5, 11, 0.8)
        tf3.text = f"ผู้นำเสนอ: {author}"
        _set_pptx_font(tf3.paragraphs[0].runs[0], font_name, 16, RGBColor(0x7F, 0x8C, 0x8D))
        tf3.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_content_slide(slide, data, title_color, accent_color, font_name):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # หัวข้อ
    title_text = data.get("title", "")
    if title_text:
        tf = _add_textbox(slide, 0.5, 0.3, 12, 1)
        tf.text = title_text
        _set_pptx_font(tf.paragraphs[0].runs[0], font_name, 28, title_color, bold=True)

    # เส้นแบ่ง
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # Bullets
    bullets = data.get("bullets", [])
    if bullets:
        tf = _add_textbox(slide, 0.8, 1.6, 11.5, 5)
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            _set_pptx_font(p.runs[0], font_name, 18, RGBColor(0x2C, 0x3E, 0x50))
            p.space_after = Pt(12)

    # ตาราง
    tbl = data.get("table")
    if tbl and tbl.get("headers") and tbl.get("rows"):
        headers = tbl["headers"]
        rows = tbl["rows"]
        left, top = Inches(0.8), Inches(4.5)
        width, height = Inches(11), Inches(2)
        table_shape = slide.shapes.add_table(
            1 + len(rows), len(headers), left, top, width, height)
        table = table_shape.table
        for c, hdr in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = hdr
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_pptx_font(r, font_name, 14, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                if c < len(headers):
                    cell = table.cell(r + 1, c)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            _set_pptx_font(run, font_name, 12, RGBColor(0x2C, 0x3E, 0x50))


def _add_section_slide(slide, data, title_color, accent_color, font_name):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    tf = _add_textbox(slide, 1, 3, 11, 1.5)
    tf.text = data.get("title", "")
    _set_pptx_font(tf.paragraphs[0].runs[0], font_name, 32, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_closing_slide(slide, data, title_color, accent_color, font_name):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    shape.line.fill.background()

    tf = _add_textbox(slide, 1, 2.5, 11, 1.5)
    tf.text = data.get("title", "ขอบคุณค่ะ")
    _set_pptx_font(tf.paragraphs[0].runs[0], font_name, 36, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    bullets = data.get("bullets", [])
    if bullets:
        tf2 = _add_textbox(slide, 2, 4, 9, 2)
        for i, b in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = b
            _set_pptx_font(p.runs[0], font_name, 18, RGBColor(0xEC, 0xF0, 0xF1))
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(10)
